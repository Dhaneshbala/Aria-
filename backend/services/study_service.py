"""Study service — generates quizzes, flashcards, mind maps, study plans."""

import json
import os
import re
import logging
from services.ollama_service import OllamaService

logger = logging.getLogger(__name__)

try:
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
except ImportError:
    MSO_AUTO_SHAPE_TYPE = None

ollama = OllamaService()

# Image cache for PPT generation (avoids re-fetching same images)
_pptx_image_cache = {}

QUIZ_LEVELS = {
    "easy": "simple recall and basic understanding, single-step problems",
    "medium": "application and some analysis, two-step problems",
    "hard": "deep analysis, multi-step reasoning, challenging edge cases",
    "exam": "timed exam simulation — comprehensive coverage, varied question types, mark-scheme style answers",
    "olympiad": "competition-level (APSMO/AMC/IMO style), highly challenging, creative lateral thinking, proof-based",
}


class StudyService:

    async def generate_quiz(
        self,
        topic: str,
        level: str = "medium",
        num_questions: int = 5,
        model: str = "qwen3:8b",
    ) -> list[dict]:
        level_desc = QUIZ_LEVELS.get(level, QUIZ_LEVELS["medium"])
        prompt = (
            f"Generate {num_questions} multiple-choice quiz questions about: {topic}\n"
            f"Difficulty: {level_desc}\n"
            f"For each question:\n"
            f"Q[N]: [question text]\n"
            f"A) [option]\nB) [option]\nC) [option]\nD) [option]\n"
            f"Correct: [letter]\n"
            f"Explanation: [brief explanation]\n\n"
            f"Make questions appropriate for a 13-year-old student."
        )
        response = await ollama.complete(model, prompt)
        return self._parse_quiz(response)

    async def generate_flashcards(
        self, topic: str, num_cards: int = 10, model: str = "qwen3:8b"
    ) -> list[dict]:
        prompt = (
            f"Generate {num_cards} flashcards for studying: {topic}\n"
            f"Format exactly:\n"
            f"FRONT: [term or question] | BACK: [definition or answer]\n"
            f"One flashcard per line. Make them clear and memorable for a 13-year-old."
        )
        response = await ollama.complete(model, prompt)
        return self._parse_flashcards(response)

    async def generate_mindmap(self, topic: str, model: str = "qwen3:8b") -> dict:
        prompt = (
            f"Create a mind map for the topic: {topic}\n"
            f"Return ONLY valid JSON in this exact format:\n"
            f'{{"center": "{topic}", "branches": [{{"label": "Branch 1", "children": ["subtopic 1", "subtopic 2"]}}, {{"label": "Branch 2", "children": ["subtopic 3", "subtopic 4"]}}]}}\n'
            f"Include 4-6 main branches, each with 2-4 children. No extra text."
        )
        response = await ollama.complete(model, prompt)
        try:
            clean = response.strip()
            if clean.startswith("```"):
                clean = re.sub(r"```(?:json)?", "", clean).strip()
            data = json.loads(clean)
            return data
        except Exception:
            return {
                "center": topic,
                "branches": [{"label": "Could not parse", "children": []}],
            }

    async def generate_study_plan(
        self, topic: str, days: int = 7, model: str = "qwen3:8b"
    ) -> list[dict]:
        prompt = (
            f"Create a {days}-day study plan for: {topic}\n"
            f"The student is 13 years old. Format as JSON array:\n"
            f'[{{"day": 1, "title": "Day title", "tasks": ["task 1", "task 2"], "time_minutes": 45}}]\n'
            f"Include realistic daily tasks and time estimates. No extra text."
        )
        response = await ollama.complete(model, prompt)
        try:
            clean = response.strip()
            if clean.startswith("```"):
                clean = re.sub(r"```(?:json)?", "", clean).strip()
            return json.loads(clean)
        except Exception:
            return [
                {"day": i + 1, "title": f"Day {i+1}", "tasks": [], "time_minutes": 45}
                for i in range(days)
            ]

    async def generate_summary(self, text: str, model: str = "qwen3:8b") -> str:
        prompt = (
            f"Summarise the following in clear, simple language for a 13-year-old student.\n"
            f"Include: key points, important vocabulary, main ideas.\n\n{text[:4000]}"
        )
        return await ollama.complete(model, prompt)

    def _parse_quiz(self, text: str) -> list[dict]:
        questions = []
        blocks = re.split(r"Q\d+[:.)]", text)
        for block in blocks[1:]:
            lines = [l.strip() for l in block.strip().split("\n") if l.strip()]
            if not lines:
                continue
            q = {"question": lines[0], "options": [], "correct": "", "explanation": ""}
            for line in lines[1:]:
                if re.match(r"[A-D][.):]\s*", line):
                    q["options"].append(re.sub(r"^[A-D][.):]\s*", "", line))
                elif re.search(r"^correct", line, re.I):
                    m = re.search(r"[A-D]", line)
                    if m:
                        q["correct"] = m.group()
                elif "explanation" in line.lower() or (q["correct"] and len(line) > 20):
                    q["explanation"] += line + " "
            if len(q["options"]) >= 2:
                questions.append(q)
        return questions[:10]

    def _parse_flashcards(self, text: str) -> list[dict]:
        cards = []
        for line in text.split("\n"):
            if "|" in line and re.search(r"(?i)front:", line):
                parts = line.split("|")
                if len(parts) >= 2:
                    front = re.sub(r"(?i)front:\s*", "", parts[0]).strip()
                    back = re.sub(r"(?i)back:\s*", "", parts[1]).strip()
                    if front and back:
                        cards.append({"front": front, "back": back})
        return cards[:20]

    async def generate_notes(
        self, topic: str, style: str = "structured", model: str = "qwen3:8b"
    ) -> str:
        style_desc = {
            "structured": "well-organised with headings, bullet points, key terms bolded",
            "cornell": "Cornell note format: main notes right, cue questions left, summary at bottom",
            "outline": "hierarchical outline with numbered sections and sub-sections",
            "simple": "simple clear bullet points easy for a 13-year-old",
        }.get(style, "structured")
        prompt = (
            f"Create comprehensive study notes about: {topic}\n"
            f"Style: {style_desc}\n"
            f"Include: key concepts, important definitions, examples, common mistakes.\n"
            f"Clear and helpful for a 13-year-old student."
        )
        return await ollama.complete(model, prompt)

    async def generate_essay_feedback(
        self, essay: str, topic: str = "", model: str = "qwen3:8b"
    ) -> str:
        prompt = (
            f"Analyse this essay and provide detailed feedback.\n"
            f"Topic: {topic or 'not specified'}\n\n"
            f"Essay:\n{essay[:4000]}\n\n"
            f"Provide feedback on:\n"
            f"1. Thesis strength and clarity\n"
            f"2. Argument structure and logical flow\n"
            f"3. Evidence quality and use of examples\n"
            f"4. Grammar, spelling, and punctuation\n"
            f"5. Vocabulary and expression\n"
            f"6. Introduction and conclusion effectiveness\n"
            f"7. Overall coherence\n\n"
            f"For each area give a rating (Needs Work/Good/Very Good/Excellent)\n"
            f"and specific suggestions with examples of how to improve.\n"
            f"End with an overall rating and top 3 improvements to make."
        )
        return await ollama.complete(model, prompt, system="You are an expert English teacher providing constructive essay feedback for a 13-year-old student.")

    async def generate_formula_reference(
        self, topic: str, model: str = "qwen3:8b"
    ) -> str:
        prompt = (
            f"Provide a comprehensive formula reference for: {topic}\n\n"
            f"For each formula include:\n"
            f"- Name of the formula\n"
            f"- The equation (use clear notation)\n"
            f"- What each variable/symbol means\n"
            f"- When to use it\n"
            f"- A worked example\n"
            f"- Common mistakes to avoid\n\n"
            f"Format clearly with headings. Suitable for a 13-year-old student."
        )
        return await ollama.complete(model, prompt, system="You are a maths/science tutor creating a clear formula reference sheet.")

    async def generate_timeline(
        self, topic: str, model: str = "qwen3:8b"
    ) -> str:
        prompt = (
            f"Create a detailed chronological timeline for: {topic}\n\n"
            f"Format as a structured list:\n"
            f"[Date/Period] — [Event]\n"
            f"  Cause: [what led to this]\n"
            f"  Significance: [why it matters]\n\n"
            f"Include 10-15 key events. Show cause-and-effect relationships.\n"
            f"Use clear dates/periods. Suitable for a 13-year-old student."
        )
        return await ollama.complete(model, prompt, system="You are a history expert creating a clear, educational timeline.")

    async def generate_exam_questions(
        self, topic: str, num: int = 10, model: str = "qwen3:8b"
    ) -> list[dict]:
        prompt = (
            f"Create {num} exam-style questions about: {topic}\n"
            f"Mix of multiple choice, short answer, extended response.\n"
            f"Include mark allocation and model answers.\n"
            f"Format: Q[N]: question\nA) B) C) D) options\nCorrect: letter\nExplanation: answer\n"
            f"Suitable for a 13-year-old."
        )
        response = await ollama.complete(model, prompt)
        return self._parse_quiz(response)

    async def generate_pptx(
        self, topic: str, slides: int = 10, model: str = "qwen3:8b"
    ) -> bytes:
        import httpx
        from pathlib import Path

        # Use faster model for PPT content generation
        pptx_model = "llama3.2:3b"

        PRESENTON_URL  = os.environ.get("PRESENTON_URL", "http://127.0.0.1:5000")
        PRESENTON_DATA = Path.home() / "presenton_data"
        PRESENTON_USER = os.environ.get("PRESENTON_USER", "")
        PRESENTON_PASS = os.environ.get("PRESENTON_PASS", "")

        try:
            async with httpx.AsyncClient(timeout=180, follow_redirects=True) as client:
                login_resp = await client.post(
                    f"{PRESENTON_URL}/api/v1/auth/login",
                    json={"username": PRESENTON_USER, "password": PRESENTON_PASS},
                )
                if login_resp.status_code != 200:
                    raise Exception(f"Presenton login failed: {login_resp.text}")

                gen_resp = await client.post(
                    f"{PRESENTON_URL}/api/v1/ppt/presentation/generate",
                    json={
                        "content": f"{topic} for a 13 year old student",
                        "n_slides": slides,
                        "template": "general",
                        "tone": "educational",
                        "verbosity": "standard",
                        "include_title_slide": True,
                        "export_as": "pptx",
                    },
                )
                if gen_resp.status_code != 200:
                    raise Exception(f"Presenton generation failed: {gen_resp.text}")

                result = gen_resp.json()
                path = result["path"]

                if path.startswith("/app_data"):
                    local_path = PRESENTON_DATA / Path(path).relative_to("/app_data")
                    if local_path.exists():
                        return local_path.read_bytes()

                full_url = path if path.startswith("http") else f"{PRESENTON_URL}{path}"
                dl_resp = await client.get(full_url)
                if dl_resp.status_code == 200 and len(dl_resp.content) > 1000:
                    return dl_resp.content

                raise Exception(f"Could not retrieve file at path: {path}")

        except Exception as e:
            logger.warning(f"[Presenton] Failed, falling back to python-pptx: {e}")
            return await self._generate_pptx_fallback(topic, slides, pptx_model)

    async def _generate_pptx_fallback(
        self, topic: str, slides: int = 10, model: str = "qwen3:8b"
    ) -> bytes:
        import io
        import re
        import requests
        import urllib.parse
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

        # ── Step 1: AI generates structured content ──────────────────────
        plan_prompt = (
            f"You are a professional presentation designer creating a Slidesgo-quality "
            f"PowerPoint for a 13-year-old student.\n\n"
            f"TOPIC: {topic}\n"
            f"NUMBER OF SLIDES: {slides}\n\n"
            f"Generate a visual theme and exactly {slides} content slides.\n\n"
            f"THEME_BG: [hex color for dark background, e.g. #0d1117]\n"
            f"THEME_ACCENT: [vivid accent hex color, e.g. #58a6ff]\n"
            f"THEME_TEXT: [light text hex color, e.g. #f0f6fc]\n"
            f"THEME_SUB: [secondary/muted hex color, e.g. #8b949e]\n"
            f"THEME_GRADIENT: [slightly lighter than BG for panels, e.g. #161b22]\n"
            f"THEME_EMOJI: [3 emojis representing this topic]\n\n"
            f"For each slide use EXACTLY this format:\n"
            f"---\n"
            f"SLIDE [N] [type: content|highlight|quote|fact]\n"
            f"TITLE: [engaging title, max 8 words]\n"
            f"BULLET1: [fascinating fact or insight, complete sentence]\n"
            f"BULLET2: [fascinating fact or insight, complete sentence]\n"
            f"BULLET3: [fascinating fact or insight, complete sentence]\n"
            f"BULLET4: [fascinating fact or insight, complete sentence]\n"
            f"IMAGE: [2-3 word search term for a relevant photo]\n"
            f"ICON: [single relevant emoji]\n"
            f"NOTE: [1-2 sentence speaker note for this slide]\n"
            f"---\n\n"
            f"SLIDE TYPES — vary them across the presentation:\n"
            f"- content: standard slide with title + 4 bullet points\n"
            f"- highlight: one big impressive fact or statistic, fewer bullets\n"
            f"- quote: a powerful quote related to the topic (use QUOTE: field)\n"
            f'- fact: a surprising "Did you know?" style slide\n\n'
            f"RULES:\n"
            f"- Use real, verified facts\n"
            f"- Be vivid, specific, and engaging for a teenager\n"
            f"- The first content slide (after title/TOC) should be a 'highlight' or 'fact' type\n"
            f"- Include at least one 'quote' type slide\n"
            f'- Make image search terms specific (e.g. "mars planet surface" not "space")\n'
        )
        response = await ollama.complete(model, plan_prompt, max_tokens=6000)

        # ── Step 2: Parse theme colours ──────────────────────────────────
        def parse_hex(text, key, default):
            m = re.search(rf'{key}:\s*#?([A-Fa-f0-9]{{6}})', text)
            if m:
                h = m.group(1)
                return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
            h = default.lstrip('#')
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

        C_BG      = parse_hex(response, 'THEME_BG',      '#0d1117')
        C_ACCENT  = parse_hex(response, 'THEME_ACCENT',  '#58a6ff')
        C_TEXT    = parse_hex(response, 'THEME_TEXT',     '#f0f6fc')
        C_SUB     = parse_hex(response, 'THEME_SUB',     '#8b949e')
        C_PANEL   = parse_hex(response, 'THEME_GRADIENT','#161b22')
        C_DARK    = RGBColor(0x01, 0x04, 0x08)
        C_WHITE   = RGBColor(0xff, 0xff, 0xff)

        emoji_match = re.search(r'THEME_EMOJI:\s*(.+)', response)
        theme_emojis = emoji_match.group(1).strip() if emoji_match else '📚✨🎓'
        # Keep only first 3 emojis
        theme_emojis = ' '.join(theme_emojis.split()[:3])

        # ── Step 3: Parse slide content ──────────────────────────────────
        def parse_slides(text):
            blocks = re.split(r'---\s*\n', text)
            parsed = []
            for block in blocks:
                slide_match = re.search(r'SLIDE\s+(\d+)(?:\s+\[type:\s*(\w+)\])?', block)
                if not slide_match:
                    continue
                slide_type = slide_match.group(2) or 'content'
                s = {
                    'title': '', 'bullets': [], 'image': '',
                    'icon': '📚', 'type': slide_type, 'note': '', 'quote': ''
                }
                for line in block.split('\n'):
                    line = line.strip()
                    if line.startswith('TITLE:'):
                        s['title'] = line.replace('TITLE:', '').strip()
                    elif re.match(r'BULLET\d:', line):
                        b = re.sub(r'BULLET\d:\s*', '', line).strip()
                        if b:
                            s['bullets'].append(b)
                    elif line.startswith('IMAGE:'):
                        s['image'] = line.replace('IMAGE:', '').strip()
                    elif line.startswith('ICON:'):
                        s['icon'] = line.replace('ICON:', '').strip()
                    elif line.startswith('NOTE:'):
                        s['note'] = line.replace('NOTE:', '').strip()
                    elif line.startswith('QUOTE:'):
                        s['quote'] = line.replace('QUOTE:', '').strip()
                if s['title']:
                    parsed.append(s)
            return parsed

        slide_data = parse_slides(response)
        if not slide_data:
            slide_data = [{
                'title': topic, 'bullets': ['Content generated by AI'],
                'image': topic, 'icon': '📚', 'type': 'content',
                'note': '', 'quote': ''
            }]

        # ── Step 4: Image fetching (multiple sources) ────────────────────
        _img_cache = {}

        def fetch_image(query, w=960, h=540):
            cache_key = f"{query}_{w}x{h}"
            if cache_key in _img_cache:
                return _img_cache[cache_key]
            if cache_key in _pptx_image_cache:
                _img_cache[cache_key] = _pptx_image_cache[cache_key]
                return _pptx_image_cache[cache_key]

            # Source 1: Unsplash JSON API (free, 50 req/hr without key)
            try:
                search_q = urllib.parse.quote(f"{query} {topic}")
                r = requests.get(
                    f"https://unsplash.com/napi/search/photos?query={search_q}&per_page=1",
                    timeout=10,
                    headers={'User-Agent': 'Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7)'}
                )
                if r.status_code == 200:
                    data = r.json()
                    results = data.get('results', [])
                    if results:
                        img_url = results[0]['urls'].get('regular') or results[0]['urls'].get('small')
                        if img_url:
                            ir = requests.get(img_url, timeout=12, headers={'User-Agent': 'Mozilla/5.0'})
                            if ir.status_code == 200 and len(ir.content) > 5000:
                                buf = io.BytesIO(ir.content)
                                _img_cache[cache_key] = buf
                                _pptx_image_cache[cache_key] = buf
                                return buf
            except Exception:
                pass

            # Source 2: Pexels (direct image URLs)
            try:
                pexels_ids = [
                    '1108099', '2562651', '3184291', '1767434', '1454360',
                    '3694711', '1567060', '414102', '208584', '87651',
                    '2473999', '1462637', '155034', '15286', '3184465',
                ]
                idx = abs(hash(query)) % len(pexels_ids)
                pid = pexels_ids[idx]
                pr = requests.get(
                    f"https://images.pexels.com/photos/{pid}/pexels-photo-{pid}.jpeg"
                    f"?auto=compress&cs=tinysrgb&w={w}",
                    timeout=10,
                    headers={'User-Agent': 'Mozilla/5.0'}
                )
                if pr.status_code == 200 and len(pr.content) > 5000:
                    buf = io.BytesIO(pr.content)
                    _img_cache[cache_key] = buf
                    return buf
            except Exception:
                pass

            # Source 3: Lorem Picsum (random high-quality photos)
            try:
                seed = abs(hash(query + topic)) % 1000
                pr = requests.get(
                    f"https://picsum.photos/seed/{seed}/{w}/{h}",
                    timeout=10, allow_redirects=True,
                    headers={'User-Agent': 'Mozilla/5.0'}
                )
                if pr.status_code == 200 and len(pr.content) > 5000:
                    buf = io.BytesIO(pr.content)
                    _img_cache[cache_key] = buf
                    return buf
            except Exception:
                pass

            return None

        # ── Step 5: Build presentation ───────────────────────────────────
        prs = Presentation()
        prs.slide_width  = Inches(13.333)   # 16:9 widescreen
        prs.slide_height = Inches(7.5)
        W = prs.slide_width
        H = prs.slide_height
        SLIDE_W = 13.333
        SLIDE_H = 7.5

        def set_bg(slide, color):
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = color

        def set_bg_gradient(slide, c1, c2):
            fill = slide.background.fill
            fill.gradient()
            fill.gradient_stops[0].color.rgb = c1
            fill.gradient_stops[0].position = 0.0
            fill.gradient_stops[1].color.rgb = c2
            fill.gradient_stops[1].position = 1.0

        def add_text(slide, text, left, top, width, height,
                     size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
                     italic=False, wrap=True, font_name='Calibri'):
            if color is None:
                color = C_TEXT
            txb = slide.shapes.add_textbox(left, top, width, height)
            tf  = txb.text_frame
            tf.word_wrap = wrap
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text           = str(text)
            run.font.size      = Pt(size)
            run.font.bold      = bold
            run.font.italic    = italic
            run.font.color.rgb = color
            run.font.name      = font_name
            return txb

        def add_multiline(slide, lines, left, top, width, height,
                          size=14, color=None, spacing=1.4, bold_first=False,
                          font_name='Calibri'):
            if color is None:
                color = C_TEXT
            txb = slide.shapes.add_textbox(left, top, width, height)
            tf  = txb.text_frame
            tf.word_wrap = True
            for idx, line in enumerate(lines):
                if idx == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.space_after = Pt(size * spacing * 0.5)
                p.space_before = Pt(2)
                run = p.add_run()
                run.text = str(line)
                run.font.size = Pt(size)
                run.font.color.rgb = color
                run.font.name = font_name
                if bold_first and idx == 0:
                    run.font.bold = True
            return txb

        def add_panel(slide, left, top, width, height, color, radius=0.0):
            shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.fill.background()
            if radius > 0:
                shape.adjustments[0] = radius
            return shape

        def add_accent_bar(slide, left, top, width=Inches(0.8), height=Inches(0.04)):
            bar = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = C_ACCENT
            bar.line.fill.background()
            return bar

        def add_image_bg(slide, img_data, left=0, top=0, width=None, height=None):
            if img_data is None:
                return
            if width is None:
                width = W
            if height is None:
                height = H
            try:
                img_data.seek(0)
                pic = slide.shapes.add_picture(img_data, left, top, width, height)
                slide.shapes._spTree.remove(pic._element)
                slide.shapes._spTree.insert(2, pic._element)
            except Exception:
                pass

        def add_notes(slide, text):
            if text:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = text

        def add_footer(slide, slide_num, total):
            # Bottom accent bar
            add_panel(slide, Inches(0), H - Inches(0.06), W, Inches(0.06), C_ACCENT)
            # Slide number
            add_text(slide, f'{slide_num}', SLIDE_W - Inches(1.2), H - Inches(0.45),
                     Inches(0.8), Inches(0.35), size=10, color=C_SUB,
                     align=PP_ALIGN.RIGHT, font_name='Calibri')
            # Branding
            add_text(slide, 'ARIA', Inches(0.4), H - Inches(0.45),
                     Inches(1.5), Inches(0.35), size=10, color=C_SUB,
                     bold=True, font_name='Calibri')

        # ── Slide 1: Title ──────────────────────────────────────────────
        ts = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(ts, C_BG)

        # Background image with overlay
        cover_img = fetch_image(f"{topic} beautiful landscape", 1200, 675)
        if cover_img:
            add_image_bg(ts, cover_img)
            # Dark overlay
            overlay = add_panel(ts, Inches(0), Inches(0), W, H, C_DARK)
            overlay.fill.fore_color.rgb = RGBColor(
                max(0, C_BG[0] - 20), max(0, C_BG[1] - 20), max(0, C_BG[2] - 20)
            )
        else:
            # Gradient background fallback
            set_bg_gradient(ts, C_BG, C_PANEL)

        # Left content panel
        add_panel(ts, Inches(0), Inches(0), Inches(6.5), H, C_DARK)

        # Accent line at top
        add_panel(ts, Inches(0), Inches(0), W, Inches(0.07), C_ACCENT)

        # Emojis
        add_text(ts, theme_emojis, Inches(0.8), Inches(1.2), Inches(5), Inches(0.8),
                 size=32, font_name='Segoe UI Emoji')

        # Title
        add_text(ts, topic.upper(), Inches(0.8), Inches(2.2), Inches(5.2), Inches(2.0),
                 size=42, bold=True, color=C_TEXT, font_name='Calibri')

        # Accent underline
        add_accent_bar(ts, Inches(0.8), Inches(4.2), Inches(2.0), Inches(0.06))

        # Subtitle
        add_text(ts, 'An ARIA Study Presentation', Inches(0.8), Inches(4.5),
                 Inches(5), Inches(0.5), size=16, color=C_SUB, italic=True)

        # Slide count
        add_text(ts, f'{len(slide_data)} slides  ·  AI Generated  ·  {topic}',
                 Inches(0.8), Inches(5.1), Inches(5), Inches(0.4),
                 size=12, color=C_SUB)

        # Bottom bar
        add_panel(ts, Inches(0), H - Inches(0.55), W, Inches(0.55), C_ACCENT)
        add_text(ts, 'ARIA — AI Study Assistant', Inches(0.6), H - Inches(0.5),
                 Inches(5), Inches(0.4), size=12, color=C_DARK, bold=True)

        add_notes(ts, f'Presentation about {topic}. Generated by ARIA AI Study Assistant.')

        # ── Slide 2: Table of Contents ──────────────────────────────────
        toc = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(toc, C_BG)
        set_bg_gradient(toc, C_BG, C_PANEL)

        add_panel(toc, Inches(0), Inches(0), W, Inches(0.07), C_ACCENT)
        add_text(toc, 'Contents', Inches(0.8), Inches(0.4), Inches(6), Inches(0.8),
                 size=36, bold=True, color=C_TEXT)
        add_accent_bar(toc, Inches(0.8), Inches(1.15), Inches(1.5))

        cols_x = [Inches(0.6), Inches(6.8)]
        for j, s in enumerate(slide_data[:10]):
            cx = cols_x[j % 2]
            cy = Inches(1.6) + (j // 2) * Inches(1.1)

            # Number circle
            num_shape = toc.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL, cx, cy + Inches(0.05),
                Inches(0.5), Inches(0.5)
            )
            num_shape.fill.solid()
            num_shape.fill.fore_color.rgb = C_ACCENT
            num_shape.line.fill.background()
            tf = num_shape.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(j + 1)
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = C_DARK

            # Title
            add_text(toc, s['title'], cx + Inches(0.7), cy + Inches(0.05),
                     Inches(5.2), Inches(0.5), size=16, color=C_TEXT)

            # Type label
            type_label = s.get('type', 'content').capitalize()
            add_text(toc, type_label, cx + Inches(0.7), cy + Inches(0.45),
                     Inches(3), Inches(0.3), size=10, color=C_SUB, italic=True)

        add_footer(toc, 2, len(slide_data) + 3)
        add_notes(toc, 'Table of contents — overview of what this presentation covers.')

        # ── Content Slides ──────────────────────────────────────────────
        for i, s in enumerate(slide_data):
            sl = prs.slides.add_slide(prs.slide_layouts[6])
            slide_type = s.get('type', 'content')
            set_bg(sl, C_BG)

            slide_num = i + 3  # +1 for title, +1 for TOC

            if slide_type == 'quote':
                # ── Quote Slide ─────────────────────────────────────
                set_bg_gradient(sl, C_BG, C_PANEL)
                add_panel(sl, Inches(0), Inches(0), W, Inches(0.07), C_ACCENT)

                # Big quotation mark
                add_text(sl, '\u201C', Inches(1), Inches(0.8), Inches(2), Inches(1.5),
                         size=120, color=C_ACCENT, font_name='Georgia')

                # Quote text
                quote_text = s.get('quote') or s['bullets'][0] if s['bullets'] else s['title']
                add_text(sl, quote_text, Inches(1.5), Inches(2.2), Inches(10), Inches(2.5),
                         size=28, italic=True, color=C_TEXT, font_name='Georgia')

                # Attribution / slide title
                add_accent_bar(sl, Inches(1.5), Inches(4.9), Inches(1.5))
                add_text(sl, f'— {s["title"]}', Inches(1.5), Inches(5.1),
                         Inches(8), Inches(0.5), size=16, color=C_SUB)

                # Background image (subtle)
                if i % 2 == 0:
                    img = fetch_image(s['image'] or topic, 800, 450)
                    if img:
                        add_image_bg(sl, img)
                        overlay = add_panel(sl, Inches(0), Inches(0), W, H, C_BG)
                        overlay.fill.fore_color.rgb = RGBColor(
                            C_BG[0], C_BG[1], C_BG[2]
                        )

            elif slide_type == 'highlight':
                # ── Highlight / Big Fact Slide ──────────────────────
                set_bg_gradient(sl, C_BG, C_PANEL)
                add_panel(sl, Inches(0), Inches(0), W, Inches(0.07), C_ACCENT)

                # Icon
                add_text(sl, s['icon'], Inches(0.8), Inches(0.5), Inches(1), Inches(1),
                         size=48, font_name='Segoe UI Emoji')

                # Title
                add_text(sl, s['title'], Inches(0.8), Inches(1.5), Inches(11), Inches(1.0),
                         size=34, bold=True, color=C_TEXT)

                add_accent_bar(sl, Inches(0.8), Inches(2.5), Inches(2.0))

                # Big fact / first bullet
                if s['bullets']:
                    add_text(sl, s['bullets'][0], Inches(0.8), Inches(3.0),
                             Inches(11), Inches(2.0), size=22, color=C_ACCENT,
                             font_name='Calibri')

                # Remaining bullets smaller
                for j, bullet in enumerate(s['bullets'][1:3]):
                    by = Inches(4.8) + Inches(j * 0.6)
                    add_text(sl, f'→  {bullet}', Inches(1.2), by,
                             Inches(10), Inches(0.55), size=14, color=C_SUB)

                # Right side image
                img = fetch_image(s['image'] or topic, 500, 560)
                if img:
                    try:
                        img.seek(0)
                        pic = sl.shapes.add_picture(img, Inches(9.5), Inches(0.5), Inches(3.5), Inches(6.5))
                        # Send to back
                        sl.shapes._spTree.remove(pic._element)
                        sl.shapes._spTree.insert(2, pic._element)
                        # Semi-transparent overlay on image area
                        overlay = add_panel(sl, Inches(9.5), Inches(0.5), Inches(3.5), Inches(6.5), C_BG)
                        overlay.fill.fore_color.rgb = RGBColor(
                            min(255, C_BG[0] + 15), min(255, C_BG[1] + 15), min(255, C_BG[2] + 15)
                        )
                    except Exception:
                        pass

            elif slide_type == 'fact':
                # ── "Did You Know?" Slide ──────────────────────────
                set_bg_gradient(sl, C_BG, C_PANEL)
                add_panel(sl, Inches(0), Inches(0), W, Inches(0.07), C_ACCENT)

                # "Did you know?" header
                add_text(sl, 'DID YOU KNOW?', Inches(0.8), Inches(0.5), Inches(6), Inches(0.7),
                         size=14, bold=True, color=C_ACCENT, font_name='Calibri')

                add_text(sl, s['title'], Inches(0.8), Inches(1.3), Inches(11), Inches(1.5),
                         size=36, bold=True, color=C_TEXT)

                add_accent_bar(sl, Inches(0.8), Inches(2.9), Inches(2.0))

                # Facts in card layout
                for j, bullet in enumerate(s['bullets'][:3]):
                    row = j // 2
                    col = j % 2
                    cx = Inches(0.8) + Inches(col * 6.0)
                    cy = Inches(3.3) + Inches(row * 1.8)

                    card = add_panel(sl, cx, cy, Inches(5.5), Inches(1.5), C_PANEL, radius=0.05)
                    # Number badge
                    badge = sl.shapes.add_shape(
                        MSO_AUTO_SHAPE_TYPE.OVAL, cx + Inches(0.2), cy + Inches(0.3),
                        Inches(0.5), Inches(0.5)
                    )
                    badge.fill.solid()
                    badge.fill.fore_color.rgb = C_ACCENT
                    badge.line.fill.background()
                    btf = badge.text_frame
                    bp = btf.paragraphs[0]
                    bp.alignment = PP_ALIGN.CENTER
                    brun = bp.add_run()
                    brun.text = str(j + 1)
                    brun.font.size = Pt(14)
                    brun.font.bold = True
                    brun.font.color.rgb = C_DARK

                    add_text(sl, bullet, cx + Inches(0.9), cy + Inches(0.25),
                             Inches(4.4), Inches(1.0), size=13, color=C_TEXT)

            else:
                # ── Standard Content Slide ──────────────────────────
                img_left = (i % 2 == 0)
                img_data = fetch_image(s['image'] or s['title'], 550, 620)

                # Image on one side
                if img_data:
                    img_x = Inches(7.5) if img_left else Inches(0)
                    try:
                        img_data.seek(0)
                        pic = sl.shapes.add_picture(img_data, img_x, Inches(0),
                                                    Inches(5.833), H)
                        sl.shapes._spTree.remove(pic._element)
                        sl.shapes._spTree.insert(2, pic._element)
                    except Exception:
                        pass

                # Content panel
                panel_x = Inches(0) if img_left else Inches(6.5)
                panel_w = Inches(7.0) if img_left else Inches(6.833)
                add_panel(sl, panel_x, Inches(0), panel_w, H, C_PANEL)

                # Top accent bar
                add_panel(sl, Inches(0), Inches(0), W, Inches(0.07), C_ACCENT)

                # Slide number badge
                num_shape = sl.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                    panel_x + Inches(0.4), Inches(0.4),
                    Inches(0.6), Inches(0.45)
                )
                num_shape.fill.solid()
                num_shape.fill.fore_color.rgb = C_ACCENT
                num_shape.line.fill.background()
                num_shape.adjustments[0] = 0.15
                tf = num_shape.text_frame
                tf.word_wrap = False
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = str(i + 1)
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = C_DARK

                # Icon + Title
                add_text(sl, s['icon'], panel_x + Inches(0.35), Inches(1.1),
                         Inches(0.8), Inches(0.8), size=32, font_name='Segoe UI Emoji')
                add_text(sl, s['title'], panel_x + Inches(1.2), Inches(1.1),
                         Inches(5.2), Inches(0.9), size=26, bold=True, color=C_TEXT)

                # Accent underline
                add_accent_bar(sl, panel_x + Inches(0.4), Inches(2.05), Inches(1.5))

                # Bullets with coloured dots
                bullet_colors = [C_ACCENT, RGBColor(0x3f, 0xb9, 0x50), C_SUB, RGBColor(0xf0, 0x8c, 0x3a)]
                for j, bullet in enumerate(s['bullets'][:4]):
                    by = Inches(2.4) + Inches(j * 1.1)
                    dot_col = bullet_colors[j % len(bullet_colors)]

                    # Colored dot
                    dot = sl.shapes.add_shape(
                        MSO_AUTO_SHAPE_TYPE.OVAL,
                        panel_x + Inches(0.4), by + Inches(0.25),
                        Inches(0.14), Inches(0.14)
                    )
                    dot.fill.solid()
                    dot.fill.fore_color.rgb = dot_col
                    dot.line.fill.background()

                    # Bullet text
                    add_text(sl, bullet, panel_x + Inches(0.7), by,
                             Inches(5.5), Inches(0.95), size=14, color=C_TEXT)

            add_footer(sl, slide_num, len(slide_data) + 3)
            add_notes(sl, s.get('note', ''))

        # ── Final Slide: Thank You ───────────────────────────────────
        es = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(es, C_BG)
        set_bg_gradient(es, C_BG, C_PANEL)

        # Background image
        end_img = fetch_image(f"{topic} beautiful", 1200, 675)
        if end_img:
            add_image_bg(es, end_img)
            overlay = add_panel(es, Inches(0), Inches(0), W, H, C_BG)

        add_panel(es, Inches(0), Inches(0), W, Inches(0.07), C_ACCENT)

        # Emojis
        add_text(es, theme_emojis, Inches(5.0), Inches(1.5), Inches(3.3), Inches(1),
                 size=48, align=PP_ALIGN.CENTER, font_name='Segoe UI Emoji')

        # Main text
        add_text(es, 'Thanks for watching!', Inches(1), Inches(2.8),
                 Inches(11.3), Inches(1.2), size=48, bold=True, color=C_TEXT,
                 align=PP_ALIGN.CENTER)

        # Accent bar
        add_accent_bar(es, Inches(5.5), Inches(4.2), Inches(2.3))

        # Subtitle
        add_text(es, f'{topic}  ·  Made with ARIA', Inches(1), Inches(4.5),
                 Inches(11.3), Inches(0.6), size=18, color=C_SUB, italic=True,
                 align=PP_ALIGN.CENTER)

        # CTA
        add_text(es, 'Ask me anything about this topic in Chat!',
                 Inches(1), Inches(5.3), Inches(11.3), Inches(0.5),
                 size=14, color=C_ACCENT, align=PP_ALIGN.CENTER)

        add_panel(es, Inches(0), H - Inches(0.55), W, Inches(0.55), C_ACCENT)
        add_text(es, 'ARIA — AI Study Assistant', Inches(0.6), H - Inches(0.5),
                 Inches(5), Inches(0.4), size=12, color=C_DARK, bold=True)

        add_notes(es, f'End of presentation on {topic}. Generated by ARIA.')

        # ── Save ─────────────────────────────────────────────────────
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        logger.info(f"Generated PPTX: {topic} ({len(slide_data)} slides, {len(buf.getvalue())} bytes)")
        return buf.read()