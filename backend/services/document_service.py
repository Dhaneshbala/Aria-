"""
Document service — full PDF/DOCX/PPTX/XLSX/TXT/ZIP processing.

Capabilities:
  • Full text extraction with page numbers preserved
  • Smart chunking for large documents (>4000 chars)
  • Quote extraction: find quotes/passages relevant to a theme
  • Semantic search within document
  • Per-page summarisation
  • Table and list extraction
"""

import asyncio
import re
import zipfile
from pathlib import Path
from typing import Optional


class DocumentService:

    # ── Public API ────────────────────────────────────────────────────────────

    async def extract_text(self, file_bytes: bytes, filename: str) -> str:
        """Extract full text with page markers."""
        ext = Path(filename).suffix.lower()
        loop = asyncio.get_running_loop()
        return await loop.run_in_executor(None, self._extract_sync, file_bytes, ext)

    async def extract_pages(self, file_bytes: bytes, filename: str) -> list[dict]:
        """
        Extract text page-by-page.
        Returns: [{"page": 1, "text": "..."}]
        """
        ext = Path(filename).suffix.lower()
        loop = asyncio.get_running_loop()
        return await loop.run_in_executor(None, self._pages_sync, file_bytes, ext)

    async def extract_quotes(
        self, file_bytes: bytes, filename: str, theme: str
    ) -> list[dict]:
        """
        Find passages/quotes from the document relevant to a theme.
        Uses keyword matching + sentence scoring.
        Returns: [{"quote": "...", "page": N, "relevance": 0.0-1.0}]
        """
        pages = await self.extract_pages(file_bytes, filename)
        loop = asyncio.get_running_loop()
        return await loop.run_in_executor(None, self._score_quotes, pages, theme)

    async def smart_chunks(
        self, file_bytes: bytes, filename: str, max_chars: int = 3000
    ) -> list[dict]:
        """
        Split document into overlapping chunks for RAG / LLM context.
        Returns: [{"chunk": N, "text": "...", "pages": "1-3"}]
        """
        pages = await self.extract_pages(file_bytes, filename)
        return self._chunk_pages(pages, max_chars)

    async def get_context_for_query(
        self, file_bytes: bytes, filename: str, query: str, max_chars: int = 6000
    ) -> str:
        """
        Return the most relevant sections of the document for a given query.
        For short docs: return full text.
        For long docs: return scored relevant sections + page numbers.
        """
        pages = await self.extract_pages(file_bytes, filename)
        full_text = "\n\n".join(f"[Page {p['page']}]\n{p['text']}" for p in pages)

        if len(full_text) <= max_chars:
            return full_text

        # Score pages by relevance to query keywords
        keywords = self._extract_keywords(query)
        scored = []
        for page in pages:
            score = self._keyword_score(page["text"], keywords)
            scored.append((score, page))

        # Always include first and last pages (intro/conclusion)
        must_include = {pages[0]["page"], pages[-1]["page"]}

        # Sort by score, take top pages within budget
        top = sorted(scored, key=lambda x: x[0], reverse=True)
        selected = []
        chars_used = 0
        for score, page in top:
            if chars_used + len(page["text"]) > max_chars:
                break
            selected.append(page)
            chars_used += len(page["text"])

        # Add must-include pages
        selected_nums = {p["page"] for p in selected}
        for page in pages:
            if page["page"] in must_include and page["page"] not in selected_nums:
                selected.append(page)

        selected.sort(key=lambda p: p["page"])

        result = f"[Document: {filename} — showing {len(selected)}/{len(pages)} most relevant pages]\n\n"
        result += "\n\n".join(f"[Page {p['page']}]\n{p['text']}" for p in selected)
        return result

    # ── Extractors ────────────────────────────────────────────────────────────

    def _extract_sync(self, data: bytes, ext: str) -> str:
        try:
            if ext == ".pdf":
                pages = self._pdf_pages(data)
                return "\n\n".join(f"[Page {p['page']}]\n{p['text']}" for p in pages) if pages else "No text found in PDF."
            elif ext in (".docx", ".doc"):
                return self._docx(data)
            elif ext in (".pptx", ".ppt"):
                return self._pptx(data)
            elif ext in (".xlsx", ".xls"):
                return self._xlsx(data)
            elif ext == ".csv":
                return data.decode("utf-8", errors="replace")
            elif ext in (".txt", ".md", ".py", ".js", ".json", ".html", ".xml"):
                return data.decode("utf-8", errors="replace")
            elif ext == ".zip":
                return self._zip(data)
            else:
                return f"File type '{ext}' not supported for text extraction."
        except Exception as e:
            return f"Could not read file: {e}"

    def _pages_sync(self, data: bytes, ext: str) -> list[dict]:
        try:
            if ext == ".pdf":
                return self._pdf_pages(data)
            else:
                text = self._extract_sync(data, ext)
                # Treat each ~500 word block as a page for non-PDFs
                words = text.split()
                pages = []
                chunk_size = 500
                for i in range(0, len(words), chunk_size):
                    pages.append({
                        "page": len(pages) + 1,
                        "text": " ".join(words[i:i + chunk_size])
                    })
                return pages or [{"page": 1, "text": text}]
        except Exception as e:
            return [{"page": 1, "text": f"Error reading file: {e}"}]

    def _pdf_pages(self, data: bytes) -> list[dict]:
        """Extract PDF with page numbers — tries pdfminer first, falls back to PyPDF2."""
        from io import BytesIO

        # Try pdfminer (better quality)
        try:
            from pdfminer.high_level import extract_pages as pm_extract_pages
            from pdfminer.layout import LTTextContainer
            pages = []
            for page_num, page_layout in enumerate(pm_extract_pages(BytesIO(data)), start=1):
                page_text = ""
                for element in page_layout:
                    if isinstance(element, LTTextContainer):
                        page_text += element.get_text()
                if page_text.strip():
                    pages.append({"page": page_num, "text": page_text.strip()})
            if pages:
                return pages
        except ImportError:
            pass
        except Exception:
            pass

        # Fallback: PyPDF2
        try:
            import PyPDF2
            reader = PyPDF2.PdfReader(BytesIO(data))
            pages = []
            for i, page in enumerate(reader.pages, start=1):
                text = page.extract_text() or ""
                if text.strip():
                    pages.append({"page": i, "text": text.strip()})
            return pages
        except Exception as e:
            return [{"page": 1, "text": f"Could not extract PDF text: {e}"}]

    def _docx(self, data: bytes) -> str:
        from io import BytesIO
        from docx import Document
        doc = Document(BytesIO(data))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                # Preserve heading levels
                if para.style.name.startswith("Heading"):
                    parts.append(f"\n## {para.text.strip()}\n")
                else:
                    parts.append(para.text.strip())
        # Also extract tables
        for table in doc.tables:
            rows = []
            for row in table.rows:
                rows.append(" | ".join(cell.text.strip() for cell in row.cells))
            if rows:
                parts.append("\n" + "\n".join(rows) + "\n")
        return "\n".join(parts)

    def _pptx(self, data: bytes) -> str:
        from io import BytesIO
        from pptx import Presentation
        prs = Presentation(BytesIO(data))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides)

    def _xlsx(self, data: bytes) -> str:
        from io import BytesIO
        import openpyxl
        wb = openpyxl.load_workbook(BytesIO(data), read_only=True, data_only=True)
        output = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            output.append(f"[Sheet: {sheet}]")
            for row in ws.iter_rows(values_only=True):
                row_str = " | ".join(str(c) if c is not None else "" for c in row)
                if row_str.strip(" |"):
                    output.append(row_str)
        return "\n".join(output)

    def _zip(self, data: bytes) -> str:
        from io import BytesIO
        texts = []
        with zipfile.ZipFile(BytesIO(data)) as zf:
            for name in zf.namelist():
                ext = Path(name).suffix.lower()
                if ext in (".pdf", ".docx", ".txt", ".pptx", ".xlsx", ".csv", ".md"):
                    try:
                        content = zf.read(name)
                        text = self._extract_sync(content, ext)
                        texts.append(f"[File: {name}]\n{text}")
                    except Exception:
                        pass
        return "\n\n---\n\n".join(texts) if texts else "No readable files in ZIP."

    # ── Quote extraction ──────────────────────────────────────────────────────

    def _score_quotes(self, pages: list[dict], theme: str) -> list[dict]:
        """
        Find sentences from the document that are relevant to the theme.
        Returns sorted list of (quote, page, relevance_score).
        """
        keywords = self._extract_keywords(theme)
        theme_words = set(theme.lower().split())
        quotes = []

        for page in pages:
            # Split into sentences
            sentences = re.split(r'(?<=[.!?])\s+', page["text"])
            for sent in sentences:
                sent = sent.strip()
                # Filter: must be quote-worthy length (15-300 words)
                word_count = len(sent.split())
                if word_count < 8 or word_count > 150:
                    continue
                # Score by keyword overlap
                score = self._keyword_score(sent, keywords)
                if score > 0:
                    quotes.append({
                        "quote": sent,
                        "page": page["page"],
                        "relevance": round(score, 3),
                        "word_count": word_count,
                    })

        # Sort by relevance, deduplicate very similar quotes
        quotes.sort(key=lambda q: q["relevance"], reverse=True)
        deduped = []
        seen_starts = set()
        for q in quotes:
            start = q["quote"][:40].lower()
            if start not in seen_starts:
                seen_starts.add(start)
                deduped.append(q)

        return deduped[:20]  # Return top 20 most relevant quotes

    def _extract_keywords(self, query: str) -> list[str]:
        """Extract meaningful keywords from a query, removing stopwords."""
        stopwords = {
            "a", "an", "the", "is", "it", "in", "on", "at", "to", "for",
            "of", "and", "or", "but", "with", "this", "that", "from", "are",
            "be", "was", "were", "can", "could", "would", "should", "what",
            "how", "why", "when", "who", "which", "my", "me", "i", "some",
            "use", "could", "find", "get", "give", "show", "tell", "about",
        }
        words = re.findall(r'\b[a-zA-Z]{3,}\b', query.lower())
        return [w for w in words if w not in stopwords]

    def _keyword_score(self, text: str, keywords: list[str]) -> float:
        """Score text by how many keywords appear, with position weighting."""
        if not keywords or not text:
            return 0.0
        text_lower = text.lower()
        hits = 0
        for kw in keywords:
            # Exact word match
            if re.search(r'\b' + re.escape(kw) + r'\b', text_lower):
                hits += 1
            # Partial match (for stemming-like behaviour)
            elif len(kw) > 5 and kw[:5] in text_lower:
                hits += 0.5
        return hits / len(keywords)

    def _chunk_pages(self, pages: list[dict], max_chars: int) -> list[dict]:
        """Group pages into overlapping chunks for LLM context windows."""
        chunks = []
        current_text = ""
        current_pages = []
        overlap_text = ""

        for page in pages:
            page_text = f"[Page {page['page']}]\n{page['text']}"
            if len(current_text) + len(page_text) > max_chars and current_text:
                # Save current chunk
                page_range = f"{current_pages[0]}-{current_pages[-1]}" if len(current_pages) > 1 else str(current_pages[0])
                chunks.append({
                    "chunk": len(chunks) + 1,
                    "text": current_text,
                    "pages": page_range,
                })
                # Overlap: keep last ~500 chars
                overlap_text = current_text[-500:] if len(current_text) > 500 else current_text
                current_text = overlap_text + "\n\n" + page_text
                current_pages = [page["page"]]
            else:
                current_text += ("\n\n" if current_text else "") + page_text
                current_pages.append(page["page"])

        if current_text:
            page_range = f"{current_pages[0]}-{current_pages[-1]}" if len(current_pages) > 1 else str(current_pages[0])
            chunks.append({
                "chunk": len(chunks) + 1,
                "text": current_text,
                "pages": page_range,
            })

        return chunks
